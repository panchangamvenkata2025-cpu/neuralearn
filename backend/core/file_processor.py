import fitz  # PyMuPDF
from pptx import Presentation
from langchain_core.documents import Document


def extract_text_from_file(file_path):
    documents = []

    # 📄 PDF
    if file_path.endswith(".pdf"):
        doc = fitz.open(file_path)
        text = ""
        for page in doc:
            text += page.get_text()

        documents.append(Document(page_content=text))

    # 📊 PPT / PPTX
    elif file_path.endswith((".ppt", ".pptx")):
        prs = Presentation(file_path)
        text = ""

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

        documents.append(Document(page_content=text))

    else:
        raise ValueError("Unsupported file type")

    return documents